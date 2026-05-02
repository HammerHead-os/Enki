#!/usr/bin/env python3
"""Horror PPT — animated spider, drifting eyes, fog, blood drips, glitch."""

import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BLACK = RGBColor(0,0,0); RED = RGBColor(0xCC,0,0); DRED = RGBColor(0x8B,0,0)
WHITE = RGBColor(0xE8,0xE8,0xE8); DIM = RGBColor(0x99,0x99,0x99)
GREEN = RGBColor(0,0xCC,0x66); AMBER = RGBColor(0xFF,0xAA,0)
G11 = RGBColor(0x11,0x11,0x11)

with open("data/ia-members.json") as f: members = json.load(f)
with open("data/ia-federal-lobbying.json") as f: lobbying = json.load(f)
with open("data/us-state-privacy-laws.json") as f: states = json.load(f)
with open("data/bill-outcomes.json") as f: bills = json.load(f)
with open("data/internet-sector-lobbying.json") as f: sector = json.load(f)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

def bg(s):
    s.background.fill.solid(); s.background.fill.fore_color.rgb = BLACK

def txt(s, l, t, w, h, text, sz=18, color=WHITE, bold=False,
        align=PP_ALIGN.LEFT, italic=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = color
    p.font.bold = bold; p.font.name = "Courier New"
    p.font.italic = italic; p.alignment = align

def multi(s, l, t, w, h, lines, sz=14, color=WHITE, align=PP_ALIGN.LEFT, bold=False):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line; p.font.size = Pt(sz); p.font.color.rgb = color
        p.font.name = "Courier New"; p.font.bold = bold; p.alignment = align
        p.space_after = Pt(sz * 0.3)

def rect(s, l, t, w, h, fill=RGBColor(0x33,0x33,0x33), line=None):
    sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line: sh.line.color.rgb = line; sh.line.width = Pt(1)
    else: sh.line.fill.background()

def pic(s, path, l, t, w, h):
    s.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def haunt(s, spider=None):
    """Add all creepy layers to a slide."""
    # Blood drip at top
    pic(s, "assets/blood_drip.gif", 0, 0, 13.333, 0.7)
    # Cobwebs
    pic(s, "assets/cobweb_top_left.png", 0, 0, 2.5, 2.5)
    pic(s, "assets/cobweb_top_right.png", 10.8, 0, 2.5, 2.5)
    # Scratches
    pic(s, "assets/scratches.png", 2, 1, 9, 5.5)
    # Eyes watching from corners
    pic(s, "assets/blinking_eyes_small.gif", 0.2, 6.7, 1.3, 0.55)
    pic(s, "assets/blinking_eyes_small.gif", 11.8, 6.7, 1.3, 0.55)
    # Spider
    if spider:
        pic(s, "assets/dangling_spider.gif", spider[0], spider[1], 0.6, 2.1)

# ── SLIDE 1: Opening ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(6.3, 0.1))
pic(s, "assets/blinking_eyes.gif", 5.2, 0.7, 2.8, 1.1)
pic(s, "assets/glitch_overlay.gif", 1, 3.3, 11, 0.9)

txt(s, 1, 1.9, 11, 1, '"Big Brother is watching...\nbut do you know who Big Brother is?"',
    sz=28, color=RED, bold=True, align=PP_ALIGN.CENTER, italic=True)
txt(s, 1.5, 3.9, 10, 1.5, "ARE YOU REALLY IN CHARGE\nOF WHAT YOU SEE?",
    sz=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
txt(s, 3, 5.8, 7, 0.5, "▓▒░  SCROLL IF YOU DARE  ░▒▓",
    sz=16, color=DRED, align=PP_ALIGN.CENTER)

# ── SLIDE 2: What is IA? ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(12.2, 0.2))

txt(s, 0.5, 0.8, 12, 0.6, "WHAT IS THE INTERNET ASSOCIATION?",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
multi(s, 2, 2.0, 9, 3.5, [
    "Founded 2012  ·  Dissolved 2021",
    "",
    "A trade group for the biggest internet companies.",
    "Status: ACTIVE until shutdown.",
    "",
    "Headquartered in Washington, D.C.",
    "Peak membership: ~40 companies.",
    "",
    "Members don't leave — unless they go bankrupt.",
    "It's too profitable to stay.",
], sz=18, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, 1, 6.0, 11, 0.4, "[placeholder: photo of IA members at inauguration]",
    sz=13, color=DIM, align=PP_ALIGN.CENTER, italic=True)

# ── SLIDE 3: Members ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(0.3, 0.1))
pic(s, "assets/blinking_eyes.gif", 5.2, 0.9, 2.8, 1.1)

txt(s, 0.5, 0.8, 12, 0.5, "WHO ARE THE MEMBERS?",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1, 1.3, 11, 0.4, "They never really leave.",
    sz=16, color=DIM, align=PP_ALIGN.CENTER, italic=True)

all_m = members["founding_members_2012"] + members["additional_members_by_2018"]
rows = []
for i in range(0, len(all_m), 3):
    rows.append("   ".join(f"▪ {all_m[j]:<18}" for j in range(i, min(i+3, len(all_m)))))
multi(s, 0.8, 2.3, 11.5, 3.5, rows, sz=14, color=WHITE)
txt(s, 1, 6.0, 11, 0.4, "[placeholder: IA members collage with scratched eyes]",
    sz=13, color=DIM, align=PP_ALIGN.CENTER, italic=True)

# ── SLIDE 4: Offices + heatmap ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(6.3, 0.05))

txt(s, 0.5, 0.8, 12, 0.5, "WHERE ARE THEY?",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.5, 1.3, 12, 0.3,
    "IA offices sit in the exact states pushing the most privacy bills.",
    sz=14, color=DIM, align=PP_ALIGN.CENTER, italic=True)

rect(s, 0.3, 1.8, 6, 4.5, G11, DRED)
txt(s, 0.5, 1.9, 5.5, 0.4, "IA OFFICE LOCATIONS", sz=16, color=AMBER, bold=True)
off = []
for st, info in states["ia_office_locations"].items():
    tag = "◉ HQ" if info["type"] == "HQ" else "○   "
    off.append(f"  {tag}  {info['city']}, {st}")
multi(s, 0.5, 2.5, 5.5, 2, off, sz=15, color=WHITE)

rect(s, 6.8, 1.8, 6.2, 4.5, G11, DRED)
txt(s, 7, 1.9, 5.8, 0.4, "PRIVACY BILLS BY STATE (2019)", sz=16, color=AMBER, bold=True)
bc = states["state_privacy_bill_counts_2019"]
offices = states["ia_office_locations"]
heat = []
for st, count in sorted(bc.items(), key=lambda x: -x[1])[:8]:
    bar = "█" * count + "░" * (25 - count)
    tag = " ◄ IA" if st in offices else ""
    heat.append(f" {st} {bar} {count}{tag}")
multi(s, 7, 2.5, 5.8, 3, heat, sz=12, color=GREEN)
txt(s, 7, 5.4, 5.8, 0.4, "Notice the overlap?", sz=16, color=RED, bold=True,
    align=PP_ALIGN.CENTER)

# ── SLIDE 5: Bar chart (GIF) ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s)

txt(s, 0.5, 0.8, 12, 0.5, "HOW MUCH ARE THEY SPENDING?",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 0.5, 1.3, 12, 0.3, "IA Federal Lobbying — Year by Year",
    sz=14, color=DIM, align=PP_ALIGN.CENTER, italic=True)
pic(s, "assets/bar_chart_lobbying.gif", 0.8, 1.8, 7.8, 4.5)
rect(s, 9, 1.8, 4, 4.5, G11, DRED)
multi(s, 9.2, 2.0, 3.6, 4, [
    "Total: $13.9M",
    "Peak: $2.4M (2017)",
    "22 lobbyists at peak",
    "",
    "Top issues:",
    "· Consumer Protection",
    "· Privacy",
    "· Telecom / Taxes",
], sz=15, color=WHITE)

# ── SLIDE 6: Pie chart (GIF) ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(12.5, 0.1))

txt(s, 0.5, 0.8, 12, 0.5, "WHO'S SPENDING THE MOST?",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1, 1.3, 11, 0.3, "Top internet sector lobbying spenders (2020)",
    sz=14, color=DIM, align=PP_ALIGN.CENTER, italic=True)
pic(s, "assets/pie_chart_spenders.gif", 0.5, 1.7, 7, 5.3)
rect(s, 8, 1.7, 5, 5.3, G11, DRED)
multi(s, 8.2, 1.9, 4.6, 4.8, [
    "9 of the top 10 spenders",
    "were IA members.",
    "",
    "Facebook: $19.7M",
    "Amazon:   $18.7M",
    "Google:    $8.9M",
    "",
    "The companies harvesting",
    "the most data spend the",
    "most to keep it legal.",
], sz=15, color=GREEN)

# ── SLIDE 7: Bills killed ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(0.2, 0.05))
pic(s, "assets/glitch_overlay.gif", 1, 1.0, 11, 0.7)

txt(s, 0.5, 0.8, 12, 0.5, "THE BILLS THEY KILLED",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
txt(s, 1, 1.5, 11, 0.4,
    "73% — When IA opposed a bill, it usually died.",
    sz=20, color=RED, bold=True, align=PP_ALIGN.CENTER)

kw = ["PRIVACY","DATA","BIOMETRIC","KNOW","GEOLOCATION",
      "INTERNET DEVICES","HOUSEHOLD","STUDENT","DIGITAL"]
pb = [b for b in bills["bills"]
      if b["ia_position"]=="OPP" and any(k in b["title"].upper() for k in kw)]
bl = []
for b in pb:
    icon = "✗ KILLED" if b["status"]=="failed" else "⚠ SURVIVED"
    bl.append(f"  {b['bill']:<10} {b['title']:<32} {b['session']}  {icon}")
multi(s, 0.5, 2.3, 12, 3.8, bl, sz=13, color=WHITE)
txt(s, 1, 6.3, 11, 0.4,
    "Right to Know. Geolocation Privacy. Data Privacy. All killed.",
    sz=14, color=DRED, align=PP_ALIGN.CENTER, italic=True)

# ── SLIDE 8: Echo chamber ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(12.7, 2.0))
pic(s, "assets/blinking_eyes.gif", 5.2, 0.7, 2.8, 1.1)

txt(s, 0.5, 0.8, 12, 0.5, "THE ECHO CHAMBER",
    sz=30, color=RED, bold=True, align=PP_ALIGN.CENTER)
multi(s, 1.5, 2.0, 10, 4, [
    "No privacy law → they collect your data freely",
    "",
    "Your data → builds your profile",
    "Your profile → feeds you more of the same",
    "More of the same → echo chamber",
    "Echo chamber → polarization",
    "Polarization → you scroll more → they profit more",
    "",
    "Americans with ideologically consistent views",
    "have DOUBLED in two decades. (Pew Research)",
    "",
    "The less you know, the easier it is",
    "to keep you scrolling — and polarized.",
], sz=18, color=WHITE, align=PP_ALIGN.CENTER)

# ── SLIDE 9: Closing ──
s = prs.slides.add_slide(prs.slide_layouts[6]); bg(s)
haunt(s, spider=(6.3, 0.05))
pic(s, "assets/blinking_eyes.gif", 1.5, 0.5, 2.8, 1.1)
pic(s, "assets/blinking_eyes.gif", 9, 0.5, 2.8, 1.1)
pic(s, "assets/blinking_eyes.gif", 5.2, 5.5, 2.8, 1.1)
pic(s, "assets/glitch_overlay.gif", 1, 2.8, 11, 0.8)
pic(s, "assets/cobweb_bottom_left.png", 0, 5, 2.5, 2.5)
pic(s, "assets/cobweb_bottom_right.png", 10.8, 5, 2.5, 2.5)

txt(s, 1, 1.7, 11, 0.8, '"Big Brother is watching..."',
    sz=28, color=RED, bold=True, align=PP_ALIGN.CENTER, italic=True)
txt(s, 1, 3.5, 11, 1.5, "NOW YOU KNOW\nWHO BIG BROTHER IS.",
    sz=46, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

rect(s, 2.5, 5.2, 8, 0.9, RGBColor(0x11,0,0), RED)
multi(s, 2.8, 5.25, 7.5, 0.8, [
    "They killed 73% of the privacy bills meant to protect you",
    "and used that data to build echo chambers.",
], sz=16, color=RED, align=PP_ALIGN.CENTER, bold=True)

txt(s, 2, 6.3, 9, 0.5,
    "▓▒░  ARE YOU STILL IN CHARGE OF WHAT YOU SEE?  ░▒▓",
    sz=18, color=DRED, align=PP_ALIGN.CENTER)

prs.save("Big_Brother_Is_Watching.pptx")
print("✓ Done — 9 slides, maximum creepiness")
